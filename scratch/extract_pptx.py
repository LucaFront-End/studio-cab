import collections
import collections.abc
from pptx import Presentation

prs = Presentation(r'd:\Workspace\Assets\Studio CAB\CONCEPTOS DE MARCA GRUPO CAB STUDIO.pptx')
print("--- CONCEPTOS DE MARCA GRUPO CAB STUDIO.pptx ---")
for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} ---")
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            print(shape.text.strip())

try:
    from pptx import Presentation
    prs = Presentation(r"d:\Workspace\Assets\Studio CAB\CONCEPTOS DE MARCA GRUPO CAB STUDIO.pptx")
    print("=== PPTX Brand Concepts ===")
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n--- Slide {slide_idx+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                print(shape.text.strip())
except Exception as e:
    print("Error reading PPTX:", e)
